#!/usr/bin/env python3
"""Fix test count 366->367 across all files."""
import re
import os
from pathlib import Path
from pptx import Presentation
ROOT = Path(__file__).resolve().parent.parent.parent

# Fix all markdown files
md_files = []
for root, dirs, files in os.walk('.'):
    dirs[:] = [d for d in dirs if d not in ['.git', 'node_modules', '__pycache__', '.claude']]
    for f in files:
        if f.endswith('.md'):
            md_files.append(os.path.join(root, f))

print("=== Fixing .md files ===")
for filepath in sorted(md_files):
    with open(filepath, 'r') as f:
        content = f.read()
    original = content
    content = content.replace('366', '367')
    if content != original:
        with open(filepath, 'w') as f:
            f.write(content)
        print(f"  ✅ {filepath}")

# Fix JSON files
json_files = []
for root, dirs, files in os.walk('.'):
    dirs[:] = [d for d in dirs if d not in ['.git', 'node_modules', '__pycache__', '.claude']]
    for f in files:
        if f.endswith('.json'):
            json_files.append(os.path.join(root, f))

print("\n=== Fixing .json files ===")
for filepath in sorted(json_files):
    with open(filepath, 'r') as f:
        content = f.read()
    original = content
    content = content.replace('366', '367')
    if content != original:
        with open(filepath, 'w') as f:
            f.write(content)
        print(f"  ✅ {filepath}")

# Fix PPT
print("\n=== Fixing Proposal_Deck.pptx ===")
prs = Presentation(ROOT / 'Proposal_Deck.pptx')
for i, slide in enumerate(prs.slides):
    changes = False
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                old = para.text
                para.text = para.text.replace('366', '367')
                if para.text != old:
                    print(f"  Slide {i+1}: '{old[:60]}' -> '{para.text[:60]}'")
                    changes = True
    if changes:
        pass
prs.save(ROOT / 'Proposal_Deck.pptx')
print("PPT saved.")

print("\nDone.")
