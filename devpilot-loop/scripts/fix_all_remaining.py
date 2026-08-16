#!/usr/bin/env python3
"""Fix ALL remaining inconsistencies across all submission files."""
import re
import os
from pathlib import Path
ROOT = Path(__file__).resolve().parent.parent

# ========== Summary of actual counts ==========
ACTUAL = {
    'evidence': 44,      # indexed evidence files
    'tests': 366,        # total test functions
    'docs': 29,          # md + json files in devpilot-loop/docs/
    'ppt_slides': 55,    # actual slide count
}

print(f"Actual counts: {ACTUAL}")

# ========== Fix README.md ==========
print("\n=== README.md ===")
with open(ROOT / 'README.md', 'r') as f:
    readme = f.read()

readme_fixes = [
    ('# 生成 54 页专业 PPT', '# 生成 55 页专业 PPT'),
    ('336 个', '366 个'),
]
for old, new in readme_fixes:
    if old in readme:
        readme = readme.replace(old, new)
        print(f"  ✅ {old[:40]} -> {new[:40]}")

with open(ROOT / 'README.md', 'w') as f:
    f.write(readme)
print("  README.md saved.")

# ========== Fix other markdown files ==========
md_files = []
for root, dirs, files in os.walk('.'):
    # Skip .git and node_modules
    dirs[:] = [d for d in dirs if d not in ['.git', 'node_modules', '__pycache__', '.claude']]
    for f in files:
        if f.endswith('.md') and not f.startswith('.'):
            md_files.append(os.path.join(root, f))

print(f"\n=== Scanning {len(md_files)} markdown files ===")
for filepath in sorted(md_files):
    with open(filepath, 'r') as f:
        content = f.read()
    original = content
    # Fix test count
    content = re.sub(r'\b336\b', '366', content)
    # Fix slide/page count
    content = re.sub(r'\b54\s*页', '55 页', content)
    content = re.sub(r'PPT 页数.*?54', 'PPT 页数（55 页）', content)
    content = re.sub(r'54 页 PPT', '55 页 PPT', content)
    if content != original:
        with open(filepath, 'w') as f:
            f.write(content)
        print(f"  ✅ {filepath}: fixed")

# ========== Fix PPTX ==========
print("\n=== Proposal_Deck.pptx ===")
from pptx import Presentation
prs = Presentation(ROOT / 'Proposal_Deck.pptx')
print(f"Slides: {len(prs.slides)}")

for slide_idx, slide in enumerate(prs.slides):
    changes = False
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                old_text = para.text
                para.text = re.sub(r'\b336\b', '366', para.text)
                para.text = re.sub(r'54\s*页', '55 页', para.text)
                if para.text != old_text:
                    changes = True
                    print(f"  Slide {slide_idx+1}: '{old_text[:60]}' -> '{para.text[:60]}'")
    if changes:
        pass  # will save at end

prs.save(ROOT / 'Proposal_Deck.pptx')
print("PPT saved.")

print("\n=== Done ===")
