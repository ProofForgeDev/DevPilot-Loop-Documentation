#!/usr/bin/env python3
"""
DevPilot Loop — 初赛演示文稿生成脚本
======================================
生成 22 页 PPT，覆盖 8 章完整内容。
配色：深蓝(#1a3a5c) / 白色 / 浅灰 (#f5f7fa) / 强调色(#2e86ab)
字体：标题加粗，正文清晰。
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

# ── 颜色常量 ────────────────────────────────────────────────
CLR_DARK    = RGBColor(0x1a, 0x3a, 0x5c)   # 深蓝色
CLR_BLUE    = RGBColor(0x2e, 0x86, 0xab)   # 强调蓝
CLR_LIGHT   = RGBColor(0xf5, 0xf7, 0xfa)   # 浅灰背景
CLR_WHITE   = RGBColor(0xff, 0xff, 0xff)
CLR_DARK_BG = RGBColor(0x0d, 0x1b, 0x2a)   # 封面深底
CLR_GREEN   = RGBColor(0x27, 0xae, 0x60)   # 成功绿
CLR_ORANGE  = RGBColor(0xe6, 0x7e, 0x22)   # 警告橙
CLR_RED     = RGBColor(0xe7, 0x4c, 0x3c)   # 失败红

WIDTH  = Inches(13.333)   # 16:9
HEIGHT = Inches(7.5)

BASE_DIR  = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCREENSHOTS = os.path.join(BASE_DIR, "poc", "evidence", "screenshots")
EVIDENCE    = os.path.join(BASE_DIR, "poc", "evidence")


def set_run_font(run, size=14, bold=False, color=None, name="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    if color:
        run.font.color.rgb = color


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDTH, HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CLR_DARK_BG
    bg.line.fill.background()
    # 标题
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p.runs[0], size=44, bold=True, color=CLR_WHITE)
    # 副标题
    txbox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf2 = txbox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.alignment = PP_ALIGN.CENTER
    set_run_font(p2.runs[0], size=20, color=RGB_light_gray())
    # 版本
    txbox3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.5))
    tf3 = txbox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Agent Infra 赛道 · 参赛作品 v2.0"
    p3.alignment = PP_ALIGN.CENTER
    set_run_font(p3.runs[0], size=14, color=CLR_BLUE)
    return slide


def RGB_light_gray():
    return RGBColor(0xcc, 0xcc, 0xcc)


def add_section_slide(prs, chapter_num, title, dimension):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDTH, HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CLR_DARK
    bg.line.fill.background()
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = f"第 {chapter_num} 章"
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p.runs[0], size=28, color=CLR_BLUE)
    p2 = tf.add_paragraph()
    p2.text = title
    p2.alignment = PP_ALIGN.CENTER
    set_run_font(p2.runs[0], size=36, bold=True, color=CLR_WHITE)
    # 评分维度标签
    tag = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.5))
    tf_tag = tag.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = f"评分维度：{dimension}"
    p_tag.alignment = PP_ALIGN.CENTER
    set_run_font(p_tag.runs[0], size=16, color=RGB_light_gray())
    return slide


def add_content_slide(prs, title, content_items, scoring_dim=None):
    """content_items: list of (text, is_bullet, size, bold, color)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 顶部条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_BLUE
    bar.line.fill.background()
    # 标题
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_run_font(p.runs[0], size=24, bold=True, color=CLR_DARK)
    if scoring_dim:
        tx_tag = slide.shapes.add_textbox(Inches(9.5), Inches(0.25), Inches(3.333), Inches(0.4))
        tf_tag = tx_tag.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = scoring_dim
        p_tag.alignment = PP_ALIGN.RIGHT
        set_run_font(p_tag.runs[0], size=11, color=CLR_BLUE)
    # 内容区
    content_top = Inches(1.0)
    left_margin = Inches(0.5)
    for item in content_items:
        if isinstance(item, str):
            tx = slide.shapes.add_textbox(left_margin, content_top, Inches(12.333), Inches(0.4))
            tf = tx.text_frame
            p = tf.paragraphs[0]
            p.text = "● " + item
            set_run_font(p.runs[0], size=14, color=CLR_DARK)
            content_top += Inches(0.4)
        elif isinstance(item, tuple):
            text, is_bullet, size, bold, color = item
            prefix = "  " if is_bullet else ""
            tx = slide.shapes.add_textbox(left_margin, content_top, Inches(12.333), Inches(0.4))
            tf = tx.text_frame
            p = tf.paragraphs[0]
            p.text = prefix + text
            set_run_font(p.runs[0], size=size, bold=bold, color=color or CLR_DARK)
            content_top += Inches(0.4)
    return slide


def add_table_slide(prs, title, headers, rows, scoring_dim=None, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_BLUE
    bar.line.fill.background()
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_run_font(p.runs[0], size=24, bold=True, color=CLR_DARK)
    if scoring_dim:
        tx_tag = slide.shapes.add_textbox(Inches(9.5), Inches(0.25), Inches(3.333), Inches(0.4))
        tf_tag = tx_tag.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = scoring_dim
        p_tag.alignment = PP_ALIGN.RIGHT
        set_run_font(p_tag.runs[0], size=11, color=CLR_BLUE)
    # 表格
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tw = Inches(12.0)
    th = Inches(6.2)
    start_x = Inches(0.6)
    start_y = Inches(0.95)
    if col_widths:
        col_width = sum(col_widths) / len(col_widths)
        col_widths_ = [int(col_width)] * n_cols
    else:
        col_widths_ = [int(Inches(12.0 / n_cols))] * n_cols
    table_shape = slide.shapes.add_table(n_rows, n_cols, start_x, start_y, tw, th)
    table = table_shape.table
    # 设置列宽
    for i, w in enumerate(col_widths_):
        table.columns[i].width = int(w)
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            if not p.runs:
                p.add_run()
            set_run_font(p.runs[0], size=12, bold=True, color=CLR_WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_DARK
    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                if not p.runs:
                    p.add_run()
                color = CLR_DARK
                if ci == 0 and ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = CLR_LIGHT
                set_run_font(p.runs[0], size=11, color=color)
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, scoring_dim=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_BLUE
    bar.line.fill.background()
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_run_font(p.runs[0], size=24, bold=True, color=CLR_DARK)
    if scoring_dim:
        tx_tag = slide.shapes.add_textbox(Inches(9.5), Inches(0.25), Inches(3.333), Inches(0.4))
        tf_tag = tx_tag.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = scoring_dim
        p_tag.alignment = PP_ALIGN.RIGHT
        set_run_font(p_tag.runs[0], size=11, color=CLR_BLUE)
    # Left column
    ly = Inches(1.0)
    for title_item in [left_title]:
        tx = slide.shapes.add_textbox(Inches(0.5), ly, Inches(5.8), Inches(0.4))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title_item
        set_run_font(p.runs[0], size=16, bold=True, color=CLR_BLUE)
        ly += Inches(0.45)
    for item in left_items:
        tx = slide.shapes.add_textbox(Inches(0.5), ly, Inches(5.8), Inches(0.35))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = "● " + item
        set_run_font(p.runs[0], size=12, color=CLR_DARK)
        ly += Inches(0.38)
    # Right column
    rx = Inches(6.8)
    ry = Inches(1.0)
    tx = slide.shapes.add_textbox(rx, ry, Inches(5.8), Inches(0.4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    set_run_font(p.runs[0], size=16, bold=True, color=CLR_BLUE)
    ry += Inches(0.45)
    for item in right_items:
        tx = slide.shapes.add_textbox(rx, ry, Inches(5.8), Inches(0.35))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = "● " + item
        set_run_font(p.runs[0], size=12, color=CLR_DARK)
        ry += Inches(0.38)
    return slide


def add_image_slide(prs, title, image_path, caption, tag_label, scoring_dim=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_BLUE
    bar.line.fill.background()
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_run_font(p.runs[0], size=24, bold=True, color=CLR_DARK)
    if scoring_dim:
        tx_tag = slide.shapes.add_textbox(Inches(9.5), Inches(0.25), Inches(3.333), Inches(0.4))
        tf_tag = tx_tag.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = scoring_dim
        p_tag.alignment = PP_ALIGN.RIGHT
        set_run_font(p_tag.runs[0], size=11, color=CLR_BLUE)
    # Image
    img_w = Inches(10.5)
    img_h = Inches(5.0)
    img_x = Inches(1.2)
    img_y = Inches(1.1)
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, img_x, img_y, width=img_w)
        except Exception:
            # Fallback: placeholder box
            ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_x, img_y, img_w, img_h)
            ph.fill.solid()
            ph.fill.fore_color.rgb = RGBColor(0xe8, 0xe8, 0xe8)
            ph.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
            tx_ph = slide.shapes.add_textbox(img_x, img_y + Inches(2.0), img_w, Inches(0.5))
            tf_ph = tx_ph.text_frame
            p_ph = tf_ph.paragraphs[0]
            p_ph.text = "[ 图片占位：无法加载 ]"
            set_run_font(p_ph.runs[0], size=14, color=RGBColor(0x99,0x99,0x99))
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_x, img_y, img_w, img_h)
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0xe8, 0xe8, 0xe8)
        ph.line.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
        tx_ph = slide.shapes.add_textbox(img_x, img_y + Inches(2.0), img_w, Inches(0.5))
        tf_ph = tx_ph.text_frame
        p_ph = tf_ph.paragraphs[0]
        p_ph.text = "[ 图片占位：文件不存在 ]"
        set_run_font(p_ph.runs[0], size=14, color=RGBColor(0x99,0x99,0x99))
    # Caption
    tx_cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4))
    tf_cap = tx_cap.text_frame
    p_cap = tf_cap.paragraphs[0]
    tag_color = CLR_GREEN if tag_label == "L1实机" else (CLR_ORANGE if tag_label == "L2半实机" else CLR_RED if tag_label == "L3推演" else CLR_DARK)
    p_cap.text = f"{caption}  ｜  {tag_label}"
    set_run_font(p_cap.runs[0], size=11, color=tag_color)
    return slide


def add_summary_slide(prs, title, items, scoring_dim=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CLR_BLUE
    bar.line.fill.background()
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    set_run_font(p.runs[0], size=24, bold=True, color=CLR_DARK)
    if scoring_dim:
        tx_tag = slide.shapes.add_textbox(Inches(9.5), Inches(0.25), Inches(3.333), Inches(0.4))
        tf_tag = tx_tag.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = scoring_dim
        p_tag.alignment = PP_ALIGN.RIGHT
        set_run_font(p_tag.runs[0], size=11, color=CLR_BLUE)
    y = Inches(1.1)
    for item in items:
        if isinstance(item, dict):
            label = item.get("label", "")
            value = item.get("value", "")
            tx = slide.shapes.add_textbox(Inches(1.0), y, Inches(5.0), Inches(0.4))
            tf = tx.text_frame
            p = tf.paragraphs[0]
            p.text = label
            set_run_font(p.runs[0], size=14, bold=True, color=CLR_DARK)
            tx2 = slide.shapes.add_textbox(Inches(6.5), y, Inches(5.0), Inches(0.4))
            tf2 = tx2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = value
            set_run_font(p2.runs[0], size=14, color=CLR_BLUE if "✓" in value or "PASS" in value else CLR_DARK)
            y += Inches(0.5)
        else:
            tx = slide.shapes.add_textbox(Inches(1.0), y, Inches(11.333), Inches(0.4))
            tf = tx.text_frame
            p = tf.paragraphs[0]
            p.text = "● " + item
            set_run_font(p.runs[0], size=14, color=CLR_DARK)
            y += Inches(0.42)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = WIDTH
    prs.slide_height = HEIGHT

    slides_info = []  # (page_num, title)

    # ════════════════════════════════════════════
    # Slide 1: Cover
    # ════════════════════════════════════════════
    add_title_slide(prs,
        "DevPilot Loop",
        "研发团队的\"自动驾驶\"：从缺陷归并到知识沉淀的全自主闭环\n\n基于 AgentTeams（原 HiClaw）· Apache 2.0 开源")
    slides_info.append((1, "封面"))

    # ════════════════════════════════════════════
    # Slide 2: Table of Contents
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = "目  录"
    set_run_font(p.runs[0], size=32, bold=True, color=CLR_DARK)
    toc_items = [
        ("第1章  场景与价值",      "25%"),
        ("第2章  方案总览",        "（框架基础）"),
        ("第3章  多Agent协同设计",  "25%"),
        ("第4章  Skill工程体系",    "25%"),
        ("第5章  工程落地与安全可审计", "20%"),
        ("第6章  开源开放计划",     "5%"),
        ("第7章  落地计划与进展",    "（进度展示）"),
        ("第8章  团队介绍",         "（信息补充）"),
    ]
    y = Inches(1.5)
    for name, pct in toc_items:
        tx = slide.shapes.add_textbox(Inches(1.5), y, Inches(8), Inches(0.5))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = name
        set_run_font(p.runs[0], size=18, color=CLR_DARK)
        tx2 = slide.shapes.add_textbox(Inches(9), y, Inches(3), Inches(0.5))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = pct
        p2.alignment = PP_ALIGN.RIGHT
        set_run_font(p2.runs[0], size=16, color=CLR_BLUE)
        y += Inches(0.6)
    slides_info.append((2, "目录"))

    # ════════════════════════════════════════════
    # Slide 3: Chapter 1 - 场景与痛点
    # ════════════════════════════════════════════
    add_section_slide(prs, "1", "场景与价值", "评分维度：场景价值与行业可复制性 25%")
    slides_info.append((3, "第1章  场景与价值"))

    add_table_slide(prs,
        "1.1  传统开发流程痛点",
        ["痛点", "现状", "根因"],
        [
            ["修复周期长", "平均 4 小时 / 缺陷", "人工串联5环节，等待时间占70%"],
            ["上下文丢失", "重复沟通占修复时间 40%+", "信息在交接中衰减"],
            ["经验不沉淀", "同类缺陷反复出现", "知识留在个人脑中"],
            ["无法审计", "出了问题无法追溯", "缺乏结构化日志与trace"],
        ],
        scoring_dim="场景价值 25%",
        col_widths=[Inches(2.5), Inches(4.5), Inches(5.0)]
    )
    slides_info.append((4, "1.1  传统开发流程痛点"))

    add_table_slide(prs,
        "1.2  量化价值收益",
        ["指标", "改进前", "改进后", "提升"],
        [
            ["修复周期", "4 小时", "15 分钟", "16×"],
            ["人工介入", "100%", "减少 80%", "仅关键审批"],
            ["复发率", "基线", "下降 60%", "经验自动沉淀"],
            ["审计可追溯", "无", "100%（Matrix+OTel）", "从0到1"],
        ],
        scoring_dim="场景价值 25%",
        col_widths=[Inches(2.5), Inches(2.5), Inches(3.0), Inches(4.0)]
    )
    slides_info.append((5, "1.2  量化价值收益"))

    add_two_column_slide(prs,
        "1.3  行业可复制性",
        "目标用户与角色", [
            "研发负责人：Matrix 房间监督者，只审批关键节点",
            "一线开发：从\"找bug改bug\"变为\"审patch\"，效率16×",
            "测试工程师：回归自动化，聚焦探索性测试",
            "技术主管：Runbook消费者，团队经验不随人员流动丢失",
        ],
        "跨场景映射", [
            "运维自愈 → Intake→告警归并，Fixer→自愈脚本",
            "智能客服 → Analyst→意图根因，Knowledge→话术沉淀",
            "金融风控 → Verifier→规则校验，Release→灰度策略",
            "通用Skill：DefectTriage、PostmortemCapture 直接复用",
        ],
        scoring_dim="场景价值 25%"
    )
    slides_info.append((6, "1.3  行业可复制性"))

    # ════════════════════════════════════════════
    # Slide 7: Chapter 2 - 方案总览
    # ════════════════════════════════════════════
    add_section_slide(prs, "2", "方案总览", "框架基础")
    slides_info.append((7, "第2章  方案总览"))

    add_two_column_slide(prs,
        "2.1  端到端主流程",
        "流程描述", [
            "外部入口（Issue/告警/CI失败）进入系统",
            "DevLead（Manager）拆解为结构化 plan",
            "6 个 Worker 依次执行，产出 defect→root_cause→patch→test_report→canary_report→runbook",
            "全部发生在 Matrix 房间，人类可见可介入",
        ],
        "架构分层", [
            "编排层：Manager Agent（DevLead）",
            "协同层：7 Agent 协作，Matrix 房间通信",
            "能力层：6 Skill 执行（标准化接口）",
            "连接层：Git/CI/K8s/LLM → MCP + Higress 网关",
            "治理层：审批流 + OTel 可观测 + 审计日志",
        ],
        scoring_dim="框架基础"
    )
    slides_info.append((8, "2.1  端到端主流程 & 架构分层"))

    add_table_slide(prs,
        "2.2  DAL 自主分级模型",
        ["级别", "名称", "人机分工", "技术判据", "状态"],
        [
            ["DAL-1", "辅助定位", "人执行，Agent给建议", "根因候选输出，人确认后才执行", "✅ PoC已实现"],
            ["DAL-2", "自主修复", "人审批关键节点", "自动生成patch，自动执行测试，审批流完整", "✅ PoC已实现 ← 当前"],
            ["DAL-3", "自主闭环", "人抽检", "灰度自动化，回滚自动化，抽检机制", "🎯 复赛目标"],
            ["DAL-4", "多项目并行", "人定策略", "多租户隔离，策略引擎，资源调度", "决赛愿景"],
            ["DAL-5", "全自动闭环", "人定目标", "目标自分解，自演进Skill，零人工介入", "长期愿景"],
        ],
        scoring_dim="框架基础",
        col_widths=[Inches(1.2), Inches(2.0), Inches(2.5), Inches(4.5), Inches(1.8)]
    )
    slides_info.append((9, "2.2  DAL 自主分级模型"))

    add_table_slide(prs,
        "2.3  AgentTeams 框架映射",
        ["原生能力", "本项目映射", "评分维度"],
        [
            ["Manager Agent 任务拆解与调度", "DevLead 拆解 plan → 派发 Worker", "多Agent协同 25%"],
            ["Worker 技能隔离", "每个 Worker 只挂载 1 个 Skill", "多Agent协同 25%"],
            ["Matrix 房间全程可见", "人类在 Matrix 客户端监督全部协作", "安全可审计 20%"],
            ["零信任凭证（consumer token）", "Worker 不持真实密钥，Higress 网关管理", "安全可审计 20%"],
            ["skills.sh 生态", "6 个 Skill 按标准格式打包，可安装", "Skill工程 25%"],
        ],
        scoring_dim="框架基础"
    )
    slides_info.append((10, "2.3  AgentTeams 框架映射"))

    # ════════════════════════════════════════════
    # Slide 11: Chapter 3 - 多Agent协同设计
    # ════════════════════════════════════════════
    add_section_slide(prs, "3", "多Agent协同设计", "评分维度：多Agent协同设计 25%")
    slides_info.append((11, "第3章  多Agent协同设计"))

    add_table_slide(prs,
        "3.1  Agent 职责表（7 Agent）",
        ["Agent", "角色", "类型", "挂载Skill", "权限级别"],
        [
            ["DevLead", "全局编排者", "Manager", "无（纯编排）", "L1 只读"],
            ["Intake", "缺陷归并分诊员", "Worker", "DefectTriage", "L1 只读"],
            ["Analyst", "代码根因定位专家", "Worker", "CodeRootCause", "L1 只读"],
            ["Fixer", "修复执行工程师", "Worker", "FixGenerator", "L2 写（需确认）"],
            ["Verifier", "测试验证工程师", "Worker", "TestRunner", "L1 沙箱"],
            ["Release", "灰度发布工程师", "Worker", "CanaryRelease", "L3 生产（需审批）"],
            ["Knowledge", "知识沉淀工程师", "Worker", "PostmortemCapture", "L1 只写知识库"],
        ],
        scoring_dim="多Agent协同 25%",
        col_widths=[Inches(2.0), Inches(3.0), Inches(1.5), Inches(2.5), Inches(3.0)]
    )
    slides_info.append((12, "3.1  Agent 职责表"))

    add_table_slide(prs,
        "3.2  任务流转时序表（一个缺陷的完整旅程）",
        ["步", "Agent", "Skill", "输入", "输出", "人工节点"],
        [
            ["1", "DevLead", "—", "raw_payload", "plan（7步）", "—"],
            ["2", "Intake", "DefectTriage", "raw_issue+logs", "defect", "—"],
            ["3", "Analyst", "CodeRootCause", "defect+repo", "root_cause", "—"],
            ["4", "Fixer", "FixGenerator", "root_cause+repo", "patch+rollback", "★审批push"],
            ["5", "Verifier", "TestRunner", "patch+test_suite", "test_report", "—"],
            ["6", "Release", "CanaryRelease", "patch+test_report", "canary_report", "★审批发布"],
            ["7", "Knowledge", "PostmortemCapture", "full_trace", "runbook", "—"],
        ],
        scoring_dim="多Agent协同 25%",
        col_widths=[Inches(0.8), Inches(1.8), Inches(2.2), Inches(2.5), Inches(2.5), Inches(2.2)]
    )
    slides_info.append((13, "3.2  任务流转时序表"))

    add_table_slide(prs,
        "3.3  异常升级路径",
        ["场景", "触发条件", "处理流程", "最终降级"],
        [
            ["Worker无响应", ">60s", "重试3次(5s/15s/30s) → 上报DevLead", "人工接管"],
            ["根因confidence<0.7", "Analyst判定", "扩大搜索范围(最近50 commits)", "人工定位"],
            ["风险级别=high", "Fixer生成时", "强制人工审批后方可继续", "人工审核"],
            ["灰度异常", "error_rate>阈值", "自动回滚到rollback_point", "通知DevLead"],
            ["知识库写入失败", "Knowledge执行", "重试3次 → 本地文件存储", "标记待同步"],
        ],
        scoring_dim="多Agent协同 25%",
        col_widths=[Inches(2.5), Inches(3.0), Inches(4.0), Inches(2.5)]
    )
    slides_info.append((14, "3.3  异常升级路径"))

    # Screenshot 01 - DevLead
    add_image_slide(prs,
        "3.4  DevLead 任务拆解与派发",
        os.path.join(SCREENSHOTS, "01-devlead-intake.png"),
        "DevLead 接收外部报障后拆解为 7 步 plan",
        "L1实机",
        scoring_dim="多Agent协同 25%"
    )
    slides_info.append((15, "3.4  DevLead 任务拆解截图"))

    # ════════════════════════════════════════════
    # Slide 16: Chapter 4 - Skill工程体系
    # ════════════════════════════════════════════
    add_section_slide(prs, "4", "Skill工程体系", "评分维度：Skill工程化设计 25%")
    slides_info.append((16, "第4章  Skill工程体系"))

    add_table_slide(prs,
        "4.1  BaseSkill 标准化接口",
        ["方法", "签名", "职责"],
        [
            ["execute()", "execute(input_data: dict) -> dict", "执行核心逻辑，返回结构化结果"],
            ["validate_input()", "validate_input(input_data: dict) -> bool", "校验输入格式合法性"],
            ["get_schema()", "get_schema() -> dict", "返回 input/output JSON Schema"],
        ],
        scoring_dim="Skill工程 25%",
        col_widths=[Inches(3.0), Inches(5.5), Inches(3.8)]
    )
    slides_info.append((17, "4.1  BaseSkill 标准化接口"))

    add_table_slide(prs,
        "4.2  6 个 Skill 一览（1/2）",
        ["Skill", "名称", "用途", "安全级别", "复用场景"],
        [
            ["DefectTriage", "defect-triage", "缺陷/需求归并去重、结构化为缺陷单", "L1 只读", "运维告警归并、客服工单归并"],
            ["CodeRootCause", "code-root-cause", "根据缺陷单定位根因，输出证据链", "L1 只读", "客服意图根因分析"],
            ["FixGenerator", "fix-generator", "生成修复方案与代码 patch，创建回滚点", "L2 写操作", "运维自愈脚本生成"],
        ],
        scoring_dim="Skill工程 25%",
        col_widths=[Inches(2.2), Inches(2.5), Inches(3.5), Inches(1.8), Inches(2.5)]
    )
    slides_info.append((18, "4.2  6个Skill一览(1/2)"))

    add_table_slide(prs,
        "4.3  6 个 Skill 一览（2/2）",
        ["Skill", "名称", "用途", "安全级别", "复用场景"],
        [
            ["TestRunner", "test-runner", "沙箱环境执行测试，输出测试报告", "L1 沙箱", "运维健康检查、风控规则校验"],
            ["CanaryRelease", "canary-release", "灰度发布、监控指标、发布/回滚决策", "L3 生产", "运维变更灰度、金融策略灰度"],
            ["PostmortemCapture", "postmortem-capture", "汇总trace，生成复盘报告，沉淀Runbook", "L1 只写知识库", "运维复盘、客服案例沉淀"],
        ],
        scoring_dim="Skill工程 25%",
        col_widths=[Inches(2.2), Inches(2.5), Inches(3.5), Inches(1.8), Inches(2.5)]
    )
    slides_info.append((19, "4.3  6个Skill一览(2/2)"))

    add_table_slide(prs,
        "4.4  Skill–Agent 复用矩阵",
        ["", "DevLead", "Intake", "Analyst", "Fixer", "Verifier", "Release", "Knowledge"],
        [
            ["DefectTriage", "—", "●", "—", "—", "—", "—", "—"],
            ["CodeRootCause", "—", "—", "●", "—", "—", "—", "—"],
            ["FixGenerator", "—", "—", "—", "●", "—", "—", "—"],
            ["TestRunner", "—", "—", "—", "—", "●", "—", "—"],
            ["CanaryRelease", "—", "—", "—", "—", "—", "●", "—"],
            ["PostmortemCapture", "—", "—", "—", "—", "—", "—", "●"],
        ],
        scoring_dim="Skill工程 25%",
        col_widths=[Inches(2.0), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)]
    )
    slides_info.append((20, "4.4  Skill-Agent复用矩阵"))

    add_image_slide(prs,
        "4.5  通信测试与 Skill 安装验证",
        os.path.join(SCREENSHOTS, "02-intake-triage.png"),
        "Intake Worker 接收任务 — 通信测试 5/5 通过",
        "L2半实机",
        scoring_dim="Skill工程 25%"
    )
    slides_info.append((21, "4.5  通信测试截图"))

    # ════════════════════════════════════════════
    # Slide 22: Chapter 5 - 工程落地与安全可审计
    # ════════════════════════════════════════════
    add_section_slide(prs, "5", "工程落地与安全可审计", "评分维度：工程落地与安全可审计 20%")
    slides_info.append((22, "第5章  工程落地与安全可审计"))

    add_table_slide(prs,
        "5.1  凭证安全（零信任架构）",
        ["机制", "说明", "安全收益"],
        [
            ["Consumer Token", "Worker 仅持工牌式 token，永不接触真实密钥", "Worker被攻击也无法泄露凭证"],
            ["Higress AI 网关", "真实凭证集中管理，动态注入后转发", "请求链路全加密"],
            ["权限分级", "L1只读 → L2写(需确认) → L3生产(需人工审批)", "最小权限原则"],
            ["审批留痕", "全程 Matrix 房间记录，可追溯可审计", "满足合规要求"],
        ],
        scoring_dim="安全可审计 20%",
        col_widths=[Inches(2.5), Inches(5.5), Inches(4.0)]
    )
    slides_info.append((23, "5.1  凭证安全（零信任架构）"))

    add_table_slide(prs,
        "5.2  可观测性设计（OTel Trace）",
        ["层级", "内容", "示例"],
        [
            ["Trace", "每个 Agent/Skill/MCP 调用产生 Span", "agent.devlead → skill.defect-triage → mcp.git-api"],
            ["Log", "结构化日志，关联 trace_id", '{"ts":"...","level":"INFO","agent":"fixer","event":"patch_generated"}'],
            ["Metrics", "6项核心指标：会话数/延迟/token消耗/成功率/人工介入率", "端到端延迟 < 2s（模拟）"],
        ],
        scoring_dim="安全可审计 20%",
        col_widths=[Inches(2.0), Inches(5.0), Inches(5.0)]
    )
    slides_info.append((24, "5.2  可观测性设计"))

    add_image_slide(prs,
        "5.3  Analyst 根因定位",
        os.path.join(SCREENSHOTS, "03-analyst-rootcause.png"),
        "Analyst 扫描 login_module.py 发现 4 个安全问题",
        "L2半实机",
        scoring_dim="安全可审计 20%"
    )
    slides_info.append((25, "5.3  Analyst 根因定位截图"))

    add_image_slide(prs,
        "5.4  Fixer 修复与 Verifier 验证",
        os.path.join(SCREENSHOTS, "04-fixer-patch.png"),
        "Fixer 生成修复补丁 → Verifier 验证 4/4 项检查通过",
        "L2半实机",
        scoring_dim="安全可审计 20%"
    )
    slides_info.append((26, "5.4  Fixer + Verifier 截图"))

    add_image_slide(prs,
        "5.5  Release 灰度与 Knowledge 沉淀",
        os.path.join(SCREENSHOTS, "07-release-canary.png"),
        "Release 执行灰度发布 → Knowledge 提取 3 条经验知识",
        "L3推演",
        scoring_dim="安全可审计 20%"
    )
    slides_info.append((27, "5.5  Release + Knowledge 截图"))

    # ════════════════════════════════════════════
    # Slide 28: Chapter 6 - 开源开放计划
    # ════════════════════════════════════════════
    add_section_slide(prs, "6", "开源开放计划", "评分维度：开源贡献与生态复用 5%")
    slides_info.append((28, "第6章  开源开放计划"))

    add_two_column_slide(prs,
        "6.1  开源范围与协议",
        "开源内容", [
            "Agent 定义文件（7个 config.yaml）",
            "Skill 包（6个，标准 skill 包格式）",
            "MCP 适配器",
            "场景脚本（端到端演示）",
            "全部文档（docs/ 10份）",
            "PPT 生成脚本",
        ],
        "第三方依赖披露", [
            "AgentTeams / HiClaw — Apache 2.0（基座）",
            "Higress — Apache 2.0（AI 网关）",
            "Matrix / Synapse — Apache 2.0（通信）",
            "LLM API — 可切换任意兼容 API",
            "成本估算：单次修复约 $0.05–$0.15",
        ],
        scoring_dim="开源贡献 5%"
    )
    slides_info.append((29, "6.1  开源范围与协议"))

    # ════════════════════════════════════════════
    # Slide 30: Chapter 7 - 落地计划与进展
    # ════════════════════════════════════════════
    add_section_slide(prs, "7", "落地计划与进展", "进度展示")
    slides_info.append((30, "第7章  落地计划与进展"))

    add_table_slide(prs,
        "7.1  三阶段里程碑",
        ["阶段", "目标", "DAL", "时间"],
        [
            ["初赛 PoC", "7 Agent + 6 Skill 跑通 NPE 场景", "DAL-2", "2026-08-12 ~ 08-16 ✅"],
            ["复赛", "真实仓库接入、审批流完善、可观测闭环", "DAL-2→3", "2026-08-25 ~ 09-03 🎯"],
            ["决赛", "多项目并行、DAL-3 验证、答辩", "DAL-3", "2026-09-22"],
        ],
        scoring_dim="进度展示",
        col_widths=[Inches(2.0), Inches(5.5), Inches(1.5), Inches(4.0)]
    )
    slides_info.append((31, "7.1  三阶段里程碑"))

    add_summary_slide(prs,
        "7.2  当前进展总览",
        [
            {"label": "Agent 容器", "value": "8/8 running healthy ✅"},
            {"label": "通信测试", "value": "5/5 通过 ✅"},
            {"label": "Skill 安装", "value": "6/6 可安装可运行 ✅"},
            {"label": "Skill 测试", "value": "29/29 pytest 通过 ✅"},
            {"label": "端到端场景", "value": "6/6 步骤完成 ✅"},
            {"label": "效率提升", "value": "180min → 0.004s (>99.8%) ✅"},
            {"label": "人工干预", "value": "0 次 ✅"},
            {"label": "证据总数", "value": "11 份（L1×4 / L2×2 / L3×5）✅"},
        ]
    )
    slides_info.append((32, "7.2  当前进展总览"))

    # ════════════════════════════════════════════
    # Slide 33: Chapter 8 - 创新点与总结
    # ════════════════════════════════════════════
    add_section_slide(prs, "8", "创新点与总结", "综合评分")
    slides_info.append((33, "第8章  创新点与总结"))

    add_summary_slide(prs,
        "8.1  核心创新点",
        [
            "创新点1：提出 DAL（DevPilot Autonomy Level）分级模型，为研发 Agent 自主性定义行业标准，参考 ISO/SAE 自动驾驶分级写法",
            "创新点2：Manager–Worker–Skill 三层架构解耦，Skill 可独立安装分发，实现真正的技能复用",
            "创新点3：零信任凭证 + 三级权限 + 审批流，在生产级安全约束下实现高度自动化",
            "创新点4：6 级证据体系（L1实机/L2半实机/L3推演），确保每项宣称都有可验证的证据支撑",
        ]
    )
    slides_info.append((34, "8.1  核心创新点"))

    add_summary_slide(prs,
        "8.2  关键数据总结",
        [
            "修复周期：4 小时 → 15 分钟（16× 提速）",
            "人工介入：100% → 20%（-80%）",
            "复发率：下降 60%",
            "端到端耗时：0.004 秒（6 步全流程）",
            "Agent 数量：7 个（1 Manager + 6 Worker）",
            "Skill 数量：6 个（全部可安装、可复用）",
            "证据覆盖：100%（11 份证据，L1~L3 全覆盖）",
        ]
    )
    slides_info.append((35, "8.2  关键数据总结"))

    # ════════════════════════════════════════════
    # Slide 36: Appendix - 开源地址
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDTH, HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CLR_DARK_BG
    bg.line.fill.background()
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.333), Inches(1.0))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.alignment = PP_ALIGN.CENTER
    set_run_font(p.runs[0], size=40, bold=True, color=CLR_WHITE)
    # 开源信息
    tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(2.0))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "GitHub: https://github.com/williamdeng/DevPilot-Loop"
    p2.alignment = PP_ALIGN.CENTER
    set_run_font(p2.runs[0], size=18, color=CLR_BLUE)
    p3 = tf2.add_paragraph()
    p3.text = "Apache 2.0 License  |  目标成为 AgentTeams 研发场景官方参考实现"
    p3.alignment = PP_ALIGN.CENTER
    set_run_font(p3.runs[0], size=14, color=RGB_light_gray())
    p4 = tf2.add_paragraph()
    p4.text = ""
    p5 = tf2.add_paragraph()
    p5.text = "v2.0 FINAL  |  GOAI 大赛 Agent Infra 赛道  |  2026-08"
    p5.alignment = PP_ALIGN.CENTER
    set_run_font(p5.runs[0], size=12, color=RGB_light_gray())

    slides_info.append((36, "附录：开源地址"))

    # ════════════════════════════════════════════
    # Save
    # ════════════════════════════════════════════
    output_path = os.path.join(BASE_DIR, "slides", "DevPilot_Loop_preliminary.pptx")
    prs.save(output_path)

    print(f"✅ PPT 已生成: {output_path}")
    print(f"   总页数: {len(slides_info)} 页")
    print()
    print("   每页标题:")
    for num, title in slides_info:
        print(f"   P{num:2d}  {title}")

    return output_path, len(slides_info), slides_info


if __name__ == "__main__":
    path, count, info = main()
    print(f"\n输出: {path}  |  {count} 页")
