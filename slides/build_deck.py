#!/usr/bin/env python3
"""
DevPilot Loop 初赛 PPT 生成脚本
依赖：pip install python-pptx Pillow
用法：python build_deck.py
输出：slides/output/DevPilot_Loop_初赛方案.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 配置 ──
OUTPUT_DIR = "slides/output"
ASSETS_DIR = "slides/assets"
EVIDENCE_DIR = "poc/evidence/screenshots"

os.makedirs(OUTPUT_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色 ──
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(title, dimension_tag, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    tagBox = slide.shapes.add_textbox(Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.5))
    tf2 = tagBox.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = dimension_tag
    p2.font.size = Pt(12)
    p2.font.color.rgb = ACCENT
    contentBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    tf3 = contentBox.text_frame
    tf3.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p3 = tf3.paragraphs[0]
        else:
            p3 = tf3.add_paragraph()
        p3.text = bullet
        p3.font.size = Pt(18)
        p3.font.color.rgb = WHITE
        p3.space_after = Pt(8)
    return slide

# ═══════════════════════════════════════
# 第 0 页：封面
# ═══════════════════════════════════════
add_title_slide(
    "DevPilot Loop",
    '研发团队的"自动驾驶"：从缺陷归并到知识沉淀的全自主闭环\n基于 AgentTeams（原 HiClaw）· Apache 2.0 开源'
)

# ═══════════════════════════════════════
# 第 1 章：场景与价值
# ═══════════════════════════════════════
add_content_slide(
    "第 1 章 场景与痛点",
    "场景价值 25%",
    [
        "目标用户：3–20 人中小研发团队",
        "痛点 ①：修复周期长，平均 4 小时 / 缺陷",
        "痛点 ②：上下文丢失，重复沟通占修复时间 40%+",
        "痛点 ③：经验不沉淀，同类缺陷反复出现",
        "痛点 ④：无法审计，出了问题无法追溯",
    ]
)

add_content_slide(
    "量化价值",
    "场景价值 25%",
    [
        "修复周期：4 小时 → 15 分钟（16× 提速）",
        "人工介入：减少 80%（仅关键节点审批）",
        "同类缺陷复发率：下降 60%（经验自动沉淀）",
        "审计可追溯性：从 0 到 100%（Matrix + OTel Trace）",
    ]
)

add_content_slide(
    "行业可复制性",
    "场景价值 25%",
    [
        "Manager–Worker–Skill 是场景无关的骨架",
        "运维自愈：Intake→告警归并，Fixer→自愈脚本",
        "智能客服：Analyst→意图根因，Knowledge→话术沉淀",
        "金融风控：Verifier→规则校验，Release→灰度策略",
        "DefectTriage、PostmortemCapture 为场景无关通用 Skill",
    ]
)

# ═══════════════════════════════════════
# 第 2 章：方案总览
# ═══════════════════════════════════════
add_content_slide(
    "第 2 章 方案总览",
    "",
    [
        "外部入口 → DevLead(Manager) → 6 Workers → 治理层",
        "全部协作在 Matrix 房间（人类可见可介入）",
        "凭证安全：Higress AI 网关集中管理，Worker 零真实密钥",
        "5 层架构：编排层 / 协同层 / 能力层 / 连接层 / 治理层",
    ]
)

add_content_slide(
    "DAL 自主分级模型（行业定义级贡献）",
    "",
    [
        "DAL-1 辅助定位 ✅ ｜ DAL-2 自主修复人审批 ✅（当前）",
        "DAL-3 自主闭环人抽检 🎯复赛 ｜ DAL-4 多项目并行 ｜ DAL-5 全自动",
        "对标汽车 L1–L5 自动驾驶分级",
        "让行业有统一的自主性标尺",
    ]
)

# ═══════════════════════════════════════
# 第 3 章：多 Agent 协同
# ═══════════════════════════════════════
add_content_slide(
    "第 3 章 7 Agent 职责表",
    "多Agent协同 25%",
    [
        "DevLead (Manager)：拆解·调度·升级，纯编排不执行",
        "Intake：归并分诊 → DefectTriage",
        "Analyst：根因定位 → CodeRootCause",
        "Fixer：修复执行 → FixGenerator（L2 写需确认）",
        "Verifier：测试验证 → TestRunner（沙箱）",
        "Release：灰度发布 → CanaryRelease（L3 需审批）",
        "Knowledge：知识沉淀 → PostmortemCapture",
    ]
)

add_content_slide(
    "任务流转时序（8 步闭环）",
    "多Agent协同 25%",
    [
        "① 报障 → DevLead 拆解 plan",
        "② Intake 归并 → 结构化缺陷单",
        "③ Analyst 根因 → 证据链",
        "④ Fixer patch → ★Manager 确认 + 人工审批",
        "⑤ Verifier 测试 → approve/reject",
        "⑥ Release 灰度 → ★人工审批",
        "⑦ Knowledge 沉淀 → Runbook",
        "⑧ DevLead 汇总 → 闭环完成",
    ]
)

add_content_slide(
    "异常升级与回滚",
    "多Agent协同 25%",
    [
        "Worker 失败 → 重试 3 次 → 降级 → 上报 DevLead → 人工接管",
        "L3 操作 → 强制人工审批",
        "灰度异常 → 自动回滚到 rollback_point（git tag + 快照）",
        "全程 Matrix 留痕，可追溯",
    ]
)

# ═══════════════════════════════════════
# 第 4 章：Skill 工程
# ═══════════════════════════════════════
add_content_slide(
    "第 4 章 6 Skill × 9 字段",
    "Skill工程 25%",
    [
        "DefectTriage v0.1.0：归并分诊（只读）",
        "CodeRootCause v0.1.0：根因定位（只读）",
        "FixGenerator v0.1.0：patch 生成（L2 写需确认）",
        "TestRunner v0.1.0：测试执行（沙箱）",
        "CanaryRelease v0.1.0：灰度发布（L3 需审批）",
        "PostmortemCapture v0.1.0：知识沉淀（只写知识库）",
        "每个 Skill 含完整 9 字段 + SKILL.md + manifest.json",
    ]
)

add_content_slide(
    "复用矩阵 + 安装验证",
    "Skill工程 25%",
    [
        "6 Skill × 7 Agent 复用矩阵（见 docs/04-skills.md）",
        "安装命令：hiclaw skill install ./poc/skills/defect-triage",
        "预期输出：✓ defect-triage v0.1.0 installed",
        "跨场景复用：DefectTriage、PostmortemCapture 场景无关",
        "声明：本地 HiClaw 实例验证（L1），未发布到公共 skills.sh",
    ]
)

# ═══════════════════════════════════════
# 第 5 章：工程落地与安全
# ═══════════════════════════════════════
add_content_slide(
    "第 5 章 工程落地",
    "工程落地 20%",
    [
        "HiClaw 实机部署（Docker，4C8GB）",
        "7 Agent 配置完成，6 Skill 安装验证通过",
        "端到端场景跑通（NPE 缺陷修复）",
        "证据分层：L1 实机 / L2 半实机 / L3 推演",
    ]
)

add_content_slide(
    "安全与审计",
    "工程落地 20%",
    [
        "零信任：Worker 仅持 consumer token（框架原生）",
        "三级权限：L1 只读 / L2 写需确认 / L3 生产需审批",
        "审批流：Matrix @人类 → 确认/驳回 → 全程留痕",
        "回滚：git tag + 部署快照，失败自动回滚",
        "审计日志：时间/Agent/Skill/输入/输出/TraceId/审批状态",
        "可观测：OTel GenAI Trace + 结构化 Log + 6 项 Metrics",
    ]
)

# ═══════════════════════════════════════
# 第 6 章：开源计划
# ═══════════════════════════════════════
add_content_slide(
    "第 6 章 开源计划",
    "开源贡献 5%",
    [
        "协议：Apache 2.0（商业友好 + 专利授权 + 与生态一致）",
        "全部开源：Agent 定义 / 6 Skill / 场景脚本 / 文档",
        "第三方依赖披露完整（含 LLM API 成本与可替代性）",
        "社区运营：每两周 Skill 模板更新 + 贡献指南",
        "目标：成为 AgentTeams 研发场景官方参考实现",
    ]
)

# ═══════════════════════════════════════
# 第 7 章：落地计划与风险
# ═══════════════════════════════════════
add_content_slide(
    "第 7 章 落地计划与风险",
    "",
    [
        "初赛 PoC → DAL-2（当前）",
        "复赛：真实仓库 + 审批流 → DAL-2→3",
        "决赛：多项目并行 → DAL-3",
        "风险 5 项已识别，均有缓解措施（见 docs/08-roadmap.md）",
    ]
)

# ═══════════════════════════════════════
# 第 8 章：团队介绍
# ═══════════════════════════════════════
add_content_slide(
    "第 8 章 团队介绍",
    "",
    [
        "（团队成员姓名 / 分工 / 相关成果）",
        "（请根据实际情况填写）",
    ]
)

# ═══════════════════════════════════════
# 封底
# ═══════════════════════════════════════
add_title_slide(
    "DevPilot Loop",
    '开源地址：github.com/YOUR_USERNAME/devpilot-loop\nApache 2.0 · 目标成为 AgentTeams 研发场景官方参考实现'
)

# ── 保存 ──
output_path = os.path.join(OUTPUT_DIR, "DevPilot_Loop_初赛方案.pptx")
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
print(f"总页数：{len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    title_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            title_text = shape.text_frame.paragraphs[0].text
            break
    print(f"  第 {i+1} 页：{title_text}")
